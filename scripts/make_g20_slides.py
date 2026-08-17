"""Generate G20_SLIDES.pptx - an easy-to-understand deck for golden case G20."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


def C(h):
    return RGBColor.from_string(h)


BG = C("F8FAFC")
DARK = C("0F172A")
MUT = C("64748B")
BLUE = C("2563EB")
GREEN = C("059669")
PURPLE = C("7C3AED")
RED = C("DC2626")
AMBER = C("D97706")
WHITE = C("FFFFFF")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def box(s, x, y, w, h, fill, txt="", size=14, color=DARK, bold=False,
        align=PP_ALIGN.LEFT, font="Segoe UI", radius=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return shp


def text(s, x, y, w, h, txt, size=14, color=DARK, bold=False,
         align=PP_ALIGN.LEFT, font="Segoe UI"):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tb


def header(s, num, title, subtitle=""):
    box(s, 0, 0, 13.333, 0.9, DARK)
    text(s, 0.5, 0.13, 9.0, 0.6, title, 26, WHITE, True)
    text(s, 12.4, 0.25, 0.7, 0.5, num, 18, C("94A3B8"), True, PP_ALIGN.RIGHT)
    if subtitle:
        text(s, 0.5, 0.95, 12.3, 0.4, subtitle, 13, MUT)


def footer(s, n):
    text(s, 0.5, 7.12, 6.0, 0.3, "G20 · Lab 17 Multi-Memory Agent", 9, MUT)
    text(s, 11.0, 7.12, 1.8, 0.3, f"{n} / 8", 9, MUT, align=PP_ALIGN.RIGHT)


def pill(s, x, y, label, color):
    box(s, x, y, 2.6, 0.42, color, label, 12, WHITE, True, PP_ALIGN.CENTER, radius=False)


# ---------- Slide 1: Title ----------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, 13.333, 7.5, BG)
text(s, 0.8, 2.1, 11.7, 1.0, "Case G20 — 3 Lớp Memory Trong 1 Câu Hỏi", 44, DARK, True, PP_ALIGN.CENTER)
text(s, 0.8, 3.05, 11.7, 0.6,
     "Short-term + Long-term + Semantic phối hợp với nhau như thế nào?", 20, MUT, False, PP_ALIGN.CENTER)
pill(s, 2.4, 4.15, "SHORT-TERM · thread", BLUE)
pill(s, 5.37, 4.15, "LONG-TERM · user", GREEN)
pill(s, 8.33, 4.15, "SEMANTIC · KB chung", PURPLE)
box(s, 4.9, 5.4, 3.55, 0.55, AMBER, "Golden 20/20  ·  +10 điểm  ·  PASS", 15, WHITE, True, PP_ALIGN.CENTER)
footer(s, 1)

# ---------- Slide 2: Bài toán ----------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, 13.333, 7.5, BG)
header(s, "2", "Bài toán của G20",
       "Một câu hỏi cần 3 nguồn thông tin khác nhau — mỗi marker đến từ một lớp")
box(s, 0.6, 1.45, 12.1, 1.5, WHITE)
text(s, 0.9, 1.6, 11.6, 0.35, "Câu hỏi (rút gọn):", 13, MUT, True)
text(s, 0.9, 1.95, 11.6, 0.95,
     "“Trong thread này mình vừa nhắc constraint giờ standup... ghép 3 mảnh: "
     "constraint còn hiệu lực trong thread, stack bắt buộc của backend công ty, và cách đánh dấu "
     "request payment để không trùng đơn.”", 16, DARK)
row = [
    ("HOLD-ALPHA-0900", "Giờ standup 09:00", "SHORT-TERM", BLUE, "Trong thread hiện tại"),
    ("NestJS", "Stack backend BLUEBIRD-42", "LONG-TERM", GREEN, "Hồ sơ của Minh (Zep)"),
    ("Idempotency-Key", "Chống trùng đơn payment", "SEMANTIC", PURPLE, "Kiến thức dùng chung"),
]
for i, (m, desc, l, c, src) in enumerate(row):
    x = 0.6 + i * 4.17
    box(s, x, 3.35, 3.9, 2.15, WHITE)
    box(s, x, 3.35, 3.9, 0.5, c, l, 12, WHITE, True, PP_ALIGN.CENTER, radius=False)
    text(s, x + 0.2, 4.0, 3.5, 0.4, m, 18, c, True, PP_ALIGN.CENTER)
    text(s, x + 0.2, 4.5, 3.5, 0.4, desc, 13, DARK, True, PP_ALIGN.CENTER)
    text(s, x + 0.2, 4.95, 3.5, 0.4, src, 11, MUT, False, PP_ALIGN.CENTER)
box(s, 0.6, 5.75, 12.1, 0.75, RED)
text(s, 0.9, 5.9, 11.6, 0.45,
     "❌ LOTUS-88 (dự án của Lan) KHÔNG được xuất hiện — kiểm tra user isolation", 15, RED, True)
footer(s, 2)

# ---------- Slide 3: Luồng ----------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, 13.333, 7.5, BG)
header(s, "3", "Luồng xử lý tổng quan",
       "Evaluator gọi 3 layer theo retrieve_layers, gộp theo budget, rồi chấm")
steps = [
    ("1. Short-term (local)", "14 tin nhắn của thread\n→ compaction → durable notes", BLUE),
    ("2. Long-term (Zep)", "User graph\n→ Context Block + facts", GREEN),
    ("3. Semantic (Zep)", "KB dùng chung\n→ document PAYMENT-RULE-3", PURPLE),
]
for i, (t, d, c) in enumerate(steps):
    x = 0.6 + i * 4.0
    box(s, x, 1.5, 3.7, 1.9, WHITE)
    box(s, x, 1.5, 3.7, 0.55, c, t, 14, WHITE, True, PP_ALIGN.CENTER, radius=False)
    for j, line in enumerate(d.split("\n")):
        text(s, x + 0.2, 2.25 + j * 0.5, 3.3, 0.5, line, 13, DARK, j == 0)
    if i < 2:
        text(s, x + 3.75, 2.1, 0.35, 0.7, "→", 24, AMBER, True, PP_ALIGN.CENTER)
box(s, 1.3, 3.9, 10.7, 0.75, AMBER)
text(s, 1.5, 4.02, 10.3, 0.5,
     "ContextBudgetManager — 10% / 4% / 3% / 3% + thứ tự ưu tiên", 15, WHITE, True, PP_ALIGN.CENTER)
text(s, 6.35, 3.55, 0.7, 0.4, "→", 20, DARK, True, PP_ALIGN.CENTER)
box(s, 1.3, 5.1, 10.7, 0.9, GREEN)
text(s, 1.5, 5.3, 10.3, 0.5,
     "✓ Scorer: HOLD-ALPHA-0900 ✓ NestJS ✓ Idempotency-Key ✓ không có LOTUS-88 → PASS",
     15, WHITE, True, PP_ALIGN.CENTER)
text(s, 0.6, 6.4, 12.1, 0.5,
     "Ghi chú: episodic = 0 vì G20 không yêu cầu — chỉ 3 layer trên được gọi.", 12, MUT)
footer(s, 3)

# ---------- Slide 4: Short-term ----------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, 13.333, 7.5, BG)
header(s, "4", "Lớp 1 — Short-term: giữ constraint dù thread dài",
       "Không gọi Zep — ShortTermMemory local (sliding window)")
box(s, 0.6, 1.5, 5.9, 3.3, WHITE)
text(s, 0.85, 1.62, 5.4, 0.35, "Đầu vào: 14 tin nhắn (fixture)", 13, BLUE, True)
text(s, 0.85, 2.0, 5.4, 1.7,
     "• 1 constraint:  “Constraint HOLD-ALPHA-0900: standup is 09:00 sharp...”\n"
     "• 1 lời xác nhận:  “Noted standup constraint.”\n"
     "• 12 tin filler (spacing, CSS, tests...)\n\n"
     "Sliding window giữ 6 turn gần nhất.\n12 tin filler bị nén/evict → 8 lần compaction.", 13, DARK)
box(s, 6.8, 1.5, 5.9, 3.3, WHITE)
text(s, 7.05, 1.62, 5.4, 0.35, "Đầu ra: durable notes giữ constraint", 13, BLUE, True)
text(s, 7.05, 2.05, 5.4, 1.5,
     "extract_durable_notes chỉ “thăng cấp” tin có dấu hiệu bền vững:\n"
     "từ khoá (constraint, deadline, todo...) hoặc MÃ VIẾT HOA.\n\n"
     "→ 2 durable notes: constraint + lời xác nhận", 13, DARK)
box(s, 7.05, 3.7, 5.4, 0.7, BLUE, "✓ HOLD-ALPHA-0900 nằm trong DURABLE_NOTES", 13, WHITE, True, PP_ALIGN.CENTER)
box(s, 0.6, 5.15, 12.1, 0.75, C("EFF6FF"))
text(s, 0.85, 5.28, 11.6, 0.5,
     "Con số: 169 / 800 token (không cần trim) · messages_kept 6 · durable_notes 2 · compactions 8",
     14, DARK, True)
text(s, 0.6, 6.15, 12.1, 0.5,
     "Bài học: compaction không phải tóm tắt văn hoa — nó ưu tiên state, decision, constraint.",
     13, MUT, True)
footer(s, 4)

# ---------- Slide 5: Long-term ----------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, 13.333, 7.5, BG)
header(s, "5", "Lớp 2 — Long-term: nhớ stack của từng dự án",
       "Zep user graph → Context Block + fact edges (kèm valid_at / invalid_at)")
box(s, 0.6, 1.5, 12.1, 2.2, WHITE)
text(s, 0.85, 1.62, 11.6, 0.35, "NestJS đến từ đâu?", 13, GREEN, True)
text(s, 0.85, 2.0, 11.6, 1.6,
     "Stage 3:  “...BLUEBIRD-42, backend bắt buộc dùng TypeScript với NestJS; không dùng Python...”\n"
     "→ Zep trích thành fact edge:  “Minh Nguyen requires NestJS for the BLUEBIRD-42 project backend.”\n"
     "→ Context Block USER_SUMMARY cũng lặp lại: “backend must use TypeScript with NestJS”\n"
     "   → có 2 nguồn, chắc chắn hơn.", 13, DARK)
box(s, 0.6, 3.95, 12.1, 0.75, C("ECFDF5"))
text(s, 0.85, 4.08, 11.6, 0.5,
     "Recency: constraint mới theo dự án (BLUEBIRD-42 → NestJS) thắng preference chung (Python) đúng scope project đó",
     14, DARK, True)
box(s, 0.6, 4.95, 12.1, 0.75, GREEN, "✓ NestJS có mặt (facts + Context Block)", 14, WHITE, True, PP_ALIGN.CENTER)
text(s, 0.6, 6.0, 12.1, 0.5,
     "Con số: raw 1403 token → trim còn 325 (giữ ĐẦU + ĐUÔI — marker có thể nằm ở cuối danh sách fact)",
     13, MUT, True)
text(s, 0.6, 6.45, 12.1, 0.5,
     "Mẹo: Zep có thể đánh dấu fact invalid_at (quirk) — nhưng marker vẫn còn nhờ nguồn thứ 2 (redundancy).",
     13, MUT, True)
footer(s, 5)

# ---------- Slide 6: Semantic ----------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, 13.333, 7.5, BG)
header(s, "6", "Lớp 3 — Semantic: kiến thức dùng chung",
       "Không thuộc riêng Minh hay Lan — search bằng graph_id, scope=“episodes”")
box(s, 0.6, 1.5, 5.9, 3.1, WHITE)
text(s, 0.85, 1.62, 5.4, 0.35, "Nguồn: data/knowledge.jsonl", 13, PURPLE, True)
text(s, 0.85, 2.0, 5.4, 2.1,
     "• kb-payment-retry  → PAYMENT-RULE-3\n"
     "• kb-async-http  → CONN-POOL-FIRST\n"
     "• kb-memory-privacy  → DELETE-VERIFY-ALL\n"
     "• kb-context-budget  → BUDGET-10-4-3-3\n\n"
     "Mỗi doc được seed 2 lần (JSON + text) → code dedupe, chỉ giữ summary.", 13, DARK)
box(s, 6.8, 1.5, 5.9, 3.1, WHITE)
text(s, 7.05, 1.62, 5.4, 0.35, "Evidence thật trả về", 13, PURPLE, True)
box(s, 7.05, 2.05, 5.4, 2.0, C("F5F3FF"))
text(s, 7.25, 2.2, 5.1, 1.7,
     "“For POST /payments, every retryable request MUST send the same Idempotency-Key...\n"
     "Marker: PAYMENT-RULE-3.”", 12, DARK)
box(s, 0.6, 4.9, 12.1, 0.75, C("F5F3FF"))
text(s, 0.85, 5.03, 11.6, 0.5,
     "Tại sao scope=“episodes”: giữ marker NGUYÊN VĂN (Idempotency-Key, PAYMENT-RULE-3). "
     "scope=“auto” chỉ trả fact, dễ mất mã literal.", 14, DARK, True)
box(s, 0.6, 5.9, 12.1, 0.75, PURPLE,
     "✓ Idempotency-Key có mặt · 97 / 240 token (không cần trim)", 14, WHITE, True, PP_ALIGN.CENTER)
footer(s, 6)

# ---------- Slide 7: Con số ----------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, 13.333, 7.5, BG)
header(s, "7", "Những con số thật (reports/golden_benchmark.json)",
       "Budget 10/4/3/3 của 8000 token → 800 / 320 / 240 / 240")
cols = ["Layer", "Limit", "Raw", "Dùng", "Ghi chú"]
rows = [
    ["Short-term", "800", "169", "169", "khớp budget, không trim"],
    ["Long-term", "320", "1403", "325", "trim giữ đầu + đuôi (+5 token nhãn)"],
    ["Episodic", "240", "0", "0", "G20 không yêu cầu layer này"],
    ["Semantic", "240", "97", "97", "dedupe JSON/text → khớp budget"],
]
x0, y0, w, h = 0.6, 1.55, 12.1, 0.55
box(s, x0, y0, w, h, DARK)
for j, cname in enumerate(cols):
    text(s, x0 + 0.2 + j * 2.42, y0 + 0.12, 2.3, 0.35, cname, 13, WHITE, True)
colors = [C("EFF6FF"), C("ECFDF5"), C("FFFBEB"), C("F5F3FF")]
for i, row in enumerate(rows):
    yy = y0 + h + i * 0.75
    box(s, x0, yy, w, 0.7, colors[i])
    for j, val in enumerate(row):
        text(s, x0 + 0.2 + j * 2.42, yy + 0.17, 2.3, 0.4, val, 13, DARK, j in (0, 3))
box(s, 0.6, 4.9, 5.9, 1.0, WHITE)
text(s, 0.85, 5.0, 5.4, 0.35, "Latency: 1641 ms", 15, DARK, True)
text(s, 0.85, 5.4, 5.4, 0.5, "get_user_context là call đắt nhất (~0.6-0.9s) + 2 graph search", 12, MUT)
box(s, 6.8, 4.9, 5.9, 1.0, WHITE)
text(s, 7.05, 5.0, 5.4, 0.35, "610 / 632 token · reduction 3.5%", 15, DARK, True)
text(s, 7.05, 5.4, 5.4, 0.5, "Thấp vì phải phủ 3 layer — cắt ít nhưng ĐÚNG hơn cắt nhiều mà SAI", 12, MUT)
text(s, 0.6, 6.2, 12.1, 0.5,
     "Đối chiếu: no-memory reduction ~82% nhưng chỉ pass 2/11 — cắt hết rẻ, nhưng sai.", 13, RED, True)
footer(s, 7)

# ---------- Slide 8: Kết luận ----------
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, 13.333, 7.5, BG)
header(s, "8", "Vì sao PASS và 3 bài học", "")
box(s, 0.6, 1.5, 12.1, 1.1, GREEN)
text(s, 0.9, 1.66, 11.6, 0.6,
     "✓ 3/3 marker có mặt (HOLD-ALPHA-0900 · NestJS · Idempotency-Key) và 0 forbidden (LOTUS-88) → PASS",
     17, WHITE, True)
lessons = [
    ("1. Scope là tường lửa",
     "Short-term đọc đúng thread · long-term chỉ user minh-lab17 · semantic chỉ graph_id dùng chung "
     "→ dữ liệu của Lan không bao giờ lọt vào."),
    ("2. Trim phải giữ cả 2 đầu",
     "Marker có thể nằm ở ĐẦU (user summary) hoặc ĐUÔI (fact cuối, marker cuối doc). "
     "Trim giữ đầu + đuôi 70/30 thay vì cắt hết đuôi."),
    ("3. Redundancy cứu case",
     "Cùng một thông tin xuất hiện ở nhiều nơi (facts + Context Block) → kể cả khi Zep đánh dấu "
     "invalid_at, marker vẫn còn để chấm."),
]
for i, (t, d) in enumerate(lessons):
    y = 2.95 + i * 1.35
    box(s, 0.6, y, 12.1, 1.2, WHITE)
    text(s, 0.9, y + 0.12, 11.6, 0.4, t, 16, DARK, True)
    text(s, 0.9, y + 0.55, 11.6, 0.6, d, 13, MUT)
text(s, 0.6, 6.6, 12.1, 0.5,
     "Tổng kết: G20 = một câu hỏi, ba nguồn nhớ, một budget chung — đúng scope + đủ marker = PASS.",
     14, AMBER, True)
footer(s, 8)

prs.save("/workspace/G20_SLIDES.pptx")
print("saved /workspace/G20_SLIDES.pptx with", len(prs.slides._sldIdLst), "slides")
